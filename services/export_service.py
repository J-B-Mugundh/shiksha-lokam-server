
from database import db
from bson import ObjectId
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import tempfile

async def export_lfa(lfa_id:str, fmt:str):
    lfa=await db.lfas.find_one({"_id":ObjectId(lfa_id)})
    tmp=tempfile.NamedTemporaryFile(delete=False,suffix=f".{fmt}")
    if fmt=="pdf":
        doc=SimpleDocTemplate(tmp.name)
        styles=getSampleStyleSheet()
        doc.build([Paragraph(lfa["name"],styles["Title"])])
    elif fmt=="docx":
        d=Document(); d.add_heading(lfa["name"]); d.save(tmp.name)
    elif fmt=="xlsx":
        wb=Workbook(); ws=wb.active; ws.append([lfa["name"]]); wb.save(tmp.name)
    elif fmt=="pptx":
        p=Presentation(); s=p.slides.add_slide(p.slide_layouts[0]); s.shapes.title.text=lfa["name"]; p.save(tmp.name)
    return tmp.name
